"""Генерация 10-слайдовой презентации сервиса «ИИ Пресс-секретарь» по ТЗ."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === Палитра (тёмная тема из ТЗ) ===
BG = RGBColor(0x0F, 0x17, 0x2A)          # #0f172a
SURFACE = RGBColor(0x1E, 0x29, 0x3B)     # #1e293b
SURFACE_2 = RGBColor(0x27, 0x34, 0x49)   # #273449
TEXT = RGBColor(0xE2, 0xE8, 0xF0)        # #e2e8f0
MUTED = RGBColor(0x94, 0xA3, 0xB8)       # #94a3b8
ACCENT = RGBColor(0x60, 0xA5, 0xFA)      # #60a5fa
ACCENT_DARK = RGBColor(0x25, 0x63, 0xEB) # #2563eb
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0x33, 0x41, 0x55)      # #334155

FONT = "Inter"
FONT_FALLBACK = "Arial"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_font(run, size, color=TEXT, bold=False, name=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name


def add_bg(slide, color=BG):
    """Заливает весь слайд цветом."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill=SURFACE, line=None, radius=False):
    """Добавляет прямоугольник (карточку)."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    if radius:
        try:
            shape.adjustments[0] = 0.06
        except Exception:
            pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, size=18, color=TEXT, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, name=FONT):
    """Добавляет текстовое поле."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold, name)
    return box


def add_paragraph(tf, text, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
                  space_after=8, bullet=False, name=FONT):
    """Добавляет абзац в существующий text_frame."""
    p = tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_font(run, size, color, bold, name)
    return p


def add_bullets(slide, x, y, w, h, items, size=18, color=TEXT, gap=10):
    """Добавляет маркированный список."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(gap)
        # Маркер
        run = p.add_run()
        run.text = "•  "
        set_font(run, size, ACCENT, True)
        # Текст
        run2 = p.add_run()
        run2.text = item
        set_font(run2, size, color)
    return box


def add_slide_number(slide, num, total=10):
    """Добавляет номер слайда в правый нижний угол."""
    add_text(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.6),
             Inches(0.9), Inches(0.4), f"{num} / {total}",
             size=12, color=MUTED, align=PP_ALIGN.RIGHT)


def add_header(slide, title, subtitle=None):
    """Добавляет заголовок слайда."""
    add_text(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9),
             title, size=34, color=TEXT, bold=True)
    if subtitle:
        add_text(slide, Inches(0.8), Inches(1.35), Inches(11.7), Inches(0.6),
                 subtitle, size=16, color=MUTED)


def add_footer_bar(slide):
    """Добавляет нижнюю полосу с брендом."""
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.35),
             fill=SURFACE_2)
    add_text(slide, Inches(0.8), SLIDE_H - Inches(0.32), Inches(6), Inches(0.3),
             "ИИ Пресс-секретарь · press.my-dpr.ru", size=11, color=MUTED)


def add_icon_circle(slide, x, y, size, emoji, bg=ACCENT):
    """Добавляет круг с эмодзи."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = emoji
    run.font.size = Pt(int(int(size) / 12700 * 0.5))
    return shape


def add_image(slide, img_path, x, y, w=None, h=None):
    """Добавляет изображение на слайд с сохранением пропорций."""
    from PIL import Image
    img = Image.open(img_path)
    iw, ih = img.size
    if w is None and h is None:
        w = Inches(4)
    if w is not None and h is None:
        w_emu = int(w)
        h = Emu(int(w_emu * ih / iw))
    elif h is not None and w is None:
        h_emu = int(h)
        w = Emu(int(h_emu * iw / ih))
    slide.shapes.add_picture(str(img_path), x, y, width=w, height=h)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # пустой макет

    # Пути к сгенерированным изображениям
    img_dir = Path(__file__).resolve().parent / "presentation_images"
    img_title = img_dir / "title_robot.png"
    img_solution = img_dir / "solution.png"
    img_media = img_dir / "media_brand.png"

    # ============ СЛАЙД 1: Титульный ============
    s = prs.slides.add_slide(blank)
    add_bg(s)
    # Изображение робота (сгенерировано через PiAPI)
    if img_title.exists():
        add_image(s, img_title, Inches(4.67), Inches(0.7), w=Inches(4.0))
    else:
        add_icon_circle(s, Inches(5.92), Inches(1.2), Inches(1.5), "🤖", ACCENT)
    # Название
    add_text(s, Inches(1.5), Inches(4.6), Inches(10.3), Inches(1.2),
             "ИИ Пресс-секретарь", size=54, color=TEXT, bold=True,
             align=PP_ALIGN.CENTER)
    # Слоган
    add_text(s, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.8),
             "Умный контент для вашего бизнеса", size=26, color=ACCENT,
             align=PP_ALIGN.CENTER)
    # Описание
    add_text(s, Inches(2.0), Inches(6.4), Inches(9.3), Inches(0.8),
             "Сервис, который создаёт контент-план, генерирует посты и картинки "
             "для вашего бизнеса и публикует их в Telegram — всё в одном месте.",
             size=15, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), SLIDE_H - Inches(0.5), Inches(6), Inches(0.3),
             "press.my-dpr.ru", size=12, color=MUTED)

    # ============ СЛАЙД 2: Проблема ============
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "Ведение соцсетей отнимает слишком много времени")
    add_footer_bar(s)
    add_slide_number(s, 2)
    problems = [
        "Нужно регулярно публиковать контент, но нет времени и идей",
        "Сложно придумывать посты каждый день",
        "Дорого нанимать копирайтера и дизайнера",
        "Трудно поддерживать единый стиль и медиа-образ",
        "Публикация в Telegram вручную — рутина",
    ]
    add_bullets(s, Inches(1.2), Inches(2.2), Inches(10.9), Inches(4.0),
                problems, size=22, gap=18)

    # ============ СЛАЙД 3: Решение ============
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "ИИ Пресс-секретарь решает эту задачу")
    add_footer_bar(s)
    add_slide_number(s, 3)
    solutions = [
        "ИИ берёт на себя генерацию контента: посты, картинки, контент-план",
        "Всё в одном личном кабинете — без копирайтера и дизайнера",
        "Контент подстраивается под ваш медиа-образ и стиль",
        "Публикация в Telegram в один клик",
        "Экономия времени и денег",
    ]
    add_bullets(s, Inches(0.8), Inches(2.2), Inches(6.5), Inches(4.0),
                solutions, size=20, gap=16)
    if img_solution.exists():
        add_image(s, img_solution, Inches(7.6), Inches(2.0), w=Inches(5.0))

    # ============ СЛАЙД 4: Ключевые возможности ============
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "Всё, что нужно для контента — в одном сервисе")
    add_footer_bar(s)
    add_slide_number(s, 4)
    features = [
        ("✍️", "Генерация постов",
         "ИИ пишет готовые тексты для вашего бизнеса с учётом вашего медиа-образа и стиля."),
        ("🎨", "Генерация картинок",
         "Создавайте иллюстрации к постам автоматически — без дизайнера и фотостока."),
        ("📅", "Контент-план",
         "Планируйте публикации на неделю вперёд и следите за выполнением плана."),
        ("📣", "Публикация в Telegram",
         "Публикуйте посты в свои каналы прямо из личного кабинета."),
    ]
    card_w = Inches(2.85)
    card_h = Inches(3.6)
    gap_x = Inches(0.25)
    start_x = Inches(0.8)
    y = Inches(2.2)
    for i, (emoji, title, desc) in enumerate(features):
        x = start_x + i * (card_w + gap_x)
        card = add_rect(s, x, y, card_w, card_h, fill=SURFACE, line=BORDER, radius=True)
        # Иконка
        add_icon_circle(s, x + Inches(0.9), y + Inches(0.4), Inches(1.0), emoji, ACCENT)
        # Заголовок
        add_text(s, x + Inches(0.3), y + Inches(1.6), card_w - Inches(0.6), Inches(0.7),
                 title, size=18, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        # Описание
        add_text(s, x + Inches(0.3), y + Inches(2.3), card_w - Inches(0.6), Inches(1.2),
                 desc, size=13, color=MUTED, align=PP_ALIGN.CENTER)

    # ============ СЛАЙД 5: Медиа-образ ============
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "Уникальный медиа-образ вашего бренда")
    add_footer_bar(s)
    add_slide_number(s, 5)
    media_items = [
        "Заполните опрос о своём бизнесе — ИИ изучит ваш стиль, миссию, ценности и целевую аудиторию",
        "Сервис формирует медиа-образ (бренд-бук) вашей организации",
        "Все посты генерируются в едином стиле, соответствующем вашему бренду",
        "Поддержка 9 видов бизнеса и различных форматов",
        "Импорт данных из существующих каналов: Telegram, VK, Tenchat, Instagram",
    ]
    add_bullets(s, Inches(0.8), Inches(2.2), Inches(6.5), Inches(4.0),
                media_items, size=18, gap=14)
    if img_media.exists():
        add_image(s, img_media, Inches(7.6), Inches(2.0), w=Inches(5.0))

    # ============ СЛАЙД 6: Как это работает ============
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "Просто начните — 4 шага")
    add_footer_bar(s)
    add_slide_number(s, 6)
    steps = [
        ("1", "Зарегистрируйтесь",
         "Создайте аккаунт в личном кабинете за пару минут."),
        ("2", "Опишите медиа-образ",
         "Расскажите о своём бизнесе — ИИ подстроится под ваш стиль."),
        ("3", "Генерируйте контент",
         "Создавайте посты, картинки и контент-план автоматически."),
        ("4", "Публикуйте",
         "Отправляйте готовые посты в свои Telegram-каналы."),
    ]
    step_w = Inches(2.85)
    step_h = Inches(3.6)
    gap_x = Inches(0.25)
    start_x = Inches(0.8)
    y = Inches(2.2)
    for i, (num, title, desc) in enumerate(steps):
        x = start_x + i * (step_w + gap_x)
        card = add_rect(s, x, y, step_w, step_h, fill=SURFACE, line=BORDER, radius=True)
        # Номер
        circle = add_rect(s, x + Inches(1.1), y + Inches(0.4), Inches(0.7), Inches(0.7),
                          fill=ACCENT, radius=True)
        tf = circle.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        set_font(run, 22, WHITE, True)
        # Заголовок
        add_text(s, x + Inches(0.3), y + Inches(1.4), step_w - Inches(0.6), Inches(0.7),
                 title, size=18, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
        # Описание
        add_text(s, x + Inches(0.3), y + Inches(2.1), step_w - Inches(0.6), Inches(1.3),
                 desc, size=13, color=MUTED, align=PP_ALIGN.CENTER)

    # ============ СЛАЙД 7: Дополнительные возможности ============
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "Больше, чем генерация постов")
    add_footer_bar(s)
    add_slide_number(s, 7)
    extra = [
        "🤖 Чат с ИИ для редактирования и доработки постов",
        "📊 Статистика и аналитика публикаций",
        "🔗 Подключение нескольких Telegram-каналов",
        "🖼️ Медиа-библиотека и история постов",
        "📱 Адаптивный дизайн — работайте с телефона, планшета или компьютера",
        "📄 Экспорт бренд-бука в документ (DOCX)",
    ]
    add_bullets(s, Inches(1.2), Inches(2.2), Inches(10.9), Inches(4.0),
                extra, size=20, gap=14)

    # ============ СЛАЙД 8: Для кого ============
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "Кому подойдёт сервис")
    add_footer_bar(s)
    add_slide_number(s, 8)
    audience = [
        "Предпринимателям и владельцам бизнеса",
        "Руководителям организаций и сообществ",
        "Маркетологам и SMM-специалистам",
        "Блогерам и экспертам, ведущим Telegram-каналы",
        "Тем, кто хочет экономить время на контенте",
    ]
    add_bullets(s, Inches(1.2), Inches(2.2), Inches(10.9), Inches(4.0),
                audience, size=22, gap=18)

    # ============ СЛАЙД 9: Преимущества ============
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "Почему выбирают ИИ Пресс-секретарь")
    add_footer_bar(s)
    add_slide_number(s, 9)
    advantages = [
        "⏱️ Экономия времени — контент создаётся за минуты",
        "💰 Экономия бюджета — не нужны копирайтер и дизайнер",
        "🎯 Единый стиль — контент соответствует вашему бренду",
        "📈 Регулярность — контент-план помогает публиковать системно",
        "🔒 Безопасность — данные и аккаунт защищены",
    ]
    add_bullets(s, Inches(1.2), Inches(2.2), Inches(10.9), Inches(4.0),
                advantages, size=22, gap=18)

    # ============ СЛАЙД 10: Призыв к действию ============
    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_icon_circle(s, Inches(6.17), Inches(1.2), Inches(1.0), "🚀", ACCENT)
    add_text(s, Inches(1.5), Inches(2.6), Inches(10.3), Inches(1.0),
             "Готовы начать?", size=44, color=TEXT, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.5), Inches(3.6), Inches(10.3), Inches(0.7),
             "Зарегистрируйтесь и получите доступ к личному кабинету.",
             size=18, color=MUTED, align=PP_ALIGN.CENTER)
    # Кнопка регистрации
    btn1 = add_rect(s, Inches(4.17), Inches(4.6), Inches(2.5), Inches(0.65),
                    fill=ACCENT, radius=True)
    tf = btn1.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Зарегистрироваться"
    set_font(run, 16, WHITE, True)
    # Кнопка входа
    btn2 = add_rect(s, Inches(6.87), Inches(4.6), Inches(2.3), Inches(0.65),
                    fill=SURFACE, line=ACCENT, radius=True)
    tf = btn2.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Войти в кабинет"
    set_font(run, 16, ACCENT, True)
    # Ссылка на сайт
    add_text(s, Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.5),
             "press.my-dpr.ru", size=16, color=ACCENT, align=PP_ALIGN.CENTER)

    # Сохранение
    out_path = Path(__file__).resolve().parent.parent / "ИИ_Пресс-секретарь_презентация.pptx"
    prs.save(out_path)
    print(f"✅ Презентация сохранена: {out_path}")


if __name__ == "__main__":
    build()
